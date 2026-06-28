#!/usr/bin/env python3
"""
project_scanner.py — monitora mudanças em arquivos do projeto e lê materiais de reunião.

Modos:
  python tools/project_scanner.py <config.json>               # delta desde o último scan
  python tools/project_scanner.py <config.json> --reset       # zera o snapshot
  python tools/project_scanner.py <config.json> --meetings    # lista materiais de reunião recentes

Saída: markdown no stdout (silencioso se sem mudanças). Atualiza o snapshot.
Arquivos abertos no Office (PermissionError) → rastreados por mtime, deep-scan adiado.
"""

import json
import os
import shutil
import sys
import tempfile
from datetime import datetime, timedelta
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import PP_PLACEHOLDER
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False


# ---------------------------------------------------------------------------
# Helpers de encoding / output
# ---------------------------------------------------------------------------

def _out(text):
    """Escreve UTF-8 no stdout sem depender do code page do console."""
    sys.stdout.buffer.write((text + "\n").encode("utf-8"))
    sys.stdout.buffer.flush()


def _fmt_ts(ts):
    return datetime.fromtimestamp(ts).strftime("%d/%m %H:%M")


def _days_ago(ts, n=3):
    return ts >= (datetime.now() - timedelta(days=n)).timestamp()


# ---------------------------------------------------------------------------
# Leitura segura de arquivos (trata arquivos abertos no Office)
# ---------------------------------------------------------------------------

def _safe_copy(path):
    """Tenta copiar o arquivo para um temp; retorna o caminho temporário ou None."""
    try:
        tmp = tempfile.NamedTemporaryFile(
            delete=False,
            suffix=Path(path).suffix,
            prefix="scanner_"
        )
        tmp.close()
        shutil.copy2(path, tmp.name)
        return tmp.name
    except Exception:
        return None


def extract_pptx_snapshot(path):
    """Extrai contagem + action titles. Tenta cópia temporária se arquivo está aberto."""
    if not HAS_PPTX:
        return {"error": "python-pptx não instalado"}
    target = path
    tmp = None
    try:
        prs = Presentation(target)
    except Exception:
        # Qualquer falha (PermissionError, PackageNotFoundError, etc.) —
        # tenta via cópia temporária (isola do lock do Office)
        tmp = _safe_copy(path)
        if tmp is None:
            return {"in_use": True, "error": "arquivo aberto ou não acessível"}
        try:
            prs = Presentation(tmp)
        except Exception as e:
            return {"in_use": True, "error": str(e)}
    finally:
        if tmp:
            try:
                os.unlink(tmp)
            except Exception:
                pass

    slides = []
    for i, slide in enumerate(prs.slides):
        slides.append({"n": i + 1, "title": _slide_title(slide)})
    return {"slide_count": len(slides), "slides": slides}


def _is_annotation_box(shape):
    """
    Heurística: shapes com fill sólido colorido (≠ branco/transparente) são caixas de
    anotação/status (ex: banner rosa de comentário, label vermelho "Preliminar"), não
    action titles. O action title BCG fica sobre fundo transparente.
    """
    try:
        from pptx.enum.dml import MSO_FILL
        fill = shape.fill
        if fill.type == MSO_FILL.SOLID:
            try:
                r, g, b = fill.fore_color.rgb
                return not (r > 230 and g > 230 and b > 230)  # qualquer cor ≠ quase-branco
            except Exception:
                return True  # theme color sólida = provavelmente annotation
    except Exception:
        pass
    return False


def _run_font_size(shape):
    """Retorna o font size em EMU do primeiro run com tamanho explícito; 0 se herdado do tema."""
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    return run.font.size
    except Exception:
        pass
    return 0


def _slide_title(slide):
    """
    Extração de action title para slides BCG/Empower:
    1. Tenta placeholder TITLE — mas só usa se NÃO existir shape não-anotação com fonte
       explicitamente maior (o Empower às vezes deixa o placeholder desatualizado).
    2. Entre shapes com texto e sem fill colorido, prefere o que tiver maior fonte
       explícita OU, se todas as fontes forem herdadas do tema, o mais alto (menor `top`).
    3. Fallback: placeholder TITLE ou qualquer shape topmost.
    """
    text_shapes = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
    non_ann = [s for s in text_shapes if not _is_annotation_box(s)]
    candidates = non_ann if non_ann else text_shapes

    # Tamanho de fonte do placeholder de título (para comparação)
    ph_title_shape = None
    ph_title_text  = None
    try:
        for ph in slide.placeholders:
            if ph.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                t = ph.text_frame.text.strip()
                if t:
                    ph_title_shape = ph
                    ph_title_text  = t
                    break
    except Exception:
        pass

    if not candidates:
        return ph_title_text[:130] if ph_title_text else "(sem título)"

    # Encontrar o shape com maior fonte explícita entre os candidatos não-anotação
    best_shape     = None
    best_font_size = 0
    for s in candidates:
        fs = _run_font_size(s)
        if fs > best_font_size:
            best_font_size = fs
            best_shape = s

    # Se existe um shape não-placeholder com fonte grande (> threshold),
    # preferi-lo sobre o placeholder (que pode estar desatualizado no Empower)
    LARGE_FONT_THRESHOLD = 200000  # ~16pt em EMU (título BCG costuma ser 24-36pt = 304800-457200 EMU)
    if best_shape and best_font_size >= LARGE_FONT_THRESHOLD and best_shape is not ph_title_shape:
        # Verificar se o placeholder também tem fonte grande — se sim, comparar textos
        ph_fs = _run_font_size(ph_title_shape) if ph_title_shape else 0
        if ph_fs >= LARGE_FONT_THRESHOLD and ph_title_text:
            # Ambos grandes: usar o que está mais acima (menor top)
            ph_top   = getattr(ph_title_shape, "top", float("inf"))
            best_top = getattr(best_shape, "top", float("inf"))
            winner = ph_title_shape if ph_top <= best_top else best_shape
            return winner.text_frame.text.strip()[:130]
        return best_shape.text_frame.text.strip()[:130]

    # Fallback: placeholder de título se disponível
    if ph_title_text:
        return ph_title_text[:130]

    # Último recurso: shape mais alto entre candidatos
    topmost = min(candidates, key=lambda s: getattr(s, "top", float("inf")))
    return topmost.text_frame.text.strip()[:130]


def extract_excel_snapshot(path, track_cells=None):
    """Extrai abas + células rastreadas. Tenta cópia se arquivo está aberto."""
    if not HAS_XLSX:
        return {"error": "openpyxl não instalado"}
    target = path
    tmp = None
    try:
        wb = openpyxl.load_workbook(target, data_only=True, read_only=True)
    except PermissionError:
        tmp = _safe_copy(path)
        if tmp is None:
            return {"in_use": True, "error": "arquivo aberto no Office"}
        try:
            wb = openpyxl.load_workbook(tmp, data_only=True, read_only=True)
        except Exception as e:
            return {"in_use": True, "error": str(e)}
    except Exception as e:
        return {"error": str(e)}
    finally:
        if tmp:
            try:
                os.unlink(tmp)
            except Exception:
                pass

    sheets = {}
    for name in wb.sheetnames:
        ws = wb[name]
        info = {"rows": ws.max_row, "cols": ws.max_column}
        if track_cells and name in track_cells:
            cells = {}
            for ref in track_cells[name]:
                try:
                    cells[ref] = str(ws[ref].value)
                except Exception:
                    cells[ref] = None
            info["cells"] = cells
        sheets[name] = info
    sheet_names = list(wb.sheetnames)
    wb.close()
    return {"sheet_names": sheet_names, "sheets": sheets}


# ---------------------------------------------------------------------------
# Scan de pasta (nomes + mtimes)
# ---------------------------------------------------------------------------

def scan_folder(folder_path):
    result = {}
    try:
        for entry in os.scandir(folder_path):
            if entry.is_file():
                result[entry.name] = round(entry.stat().st_mtime, 1)
    except Exception:
        pass
    return result


# ---------------------------------------------------------------------------
# Snapshot I/O
# ---------------------------------------------------------------------------

def load_snapshot(path):
    try:
        with open(path, encoding="utf-8") as f:
            return json.load(f)
    except (FileNotFoundError, json.JSONDecodeError):
        return {}


def save_snapshot(path, snapshot):
    Path(path).parent.mkdir(parents=True, exist_ok=True)
    with open(path, "w", encoding="utf-8") as f:
        json.dump(snapshot, f, ensure_ascii=False, indent=2)


# ---------------------------------------------------------------------------
# Delta report (modo padrão)
# ---------------------------------------------------------------------------

def build_delta_report(config, prev, curr):
    lines = []
    prev_ts = prev.get("timestamp", "—")

    # Pastas vigiadas
    for folder_cfg in config.get("watch_folders", []):
        label = folder_cfg["label"]
        path  = folder_cfg["path"]
        prev_f = prev.get("folders", {}).get(path, {})
        curr_f = curr["folders"].get(path, {})

        added    = [f for f in curr_f if f not in prev_f]
        removed  = [f for f in prev_f if f not in curr_f]
        modified = [f for f in curr_f if f in prev_f and curr_f[f] != prev_f[f]]

        if added or removed or modified:
            lines.append(f"**{label}**")
            for f in added:
                lines.append(f"  + `{f}` — novo ({_fmt_ts(curr_f[f])})")
            for f in modified:
                lines.append(f"  ~ `{f}` — modificado ({_fmt_ts(curr_f[f])})")
            for f in removed:
                lines.append(f"  - `{f}` — removido")
            lines.append("")

    # Arquivos com deep-scan
    for file_cfg in config.get("deep_scan_files", []):
        label = file_cfg["label"]
        path  = file_cfg["path"]
        ftype = file_cfg.get("type", "unknown")

        curr_data = curr["files"].get(path)
        if curr_data is None:
            continue

        prev_data  = prev.get("files", {}).get(path, {})
        prev_mtime = prev_data.get("mtime")
        curr_mtime = curr_data.get("mtime")

        if prev_mtime == curr_mtime and prev_data:
            continue  # sem mudança

        lines.append(f"**{label}**")
        lines.append(f"  `{Path(path).name}` — modificado em {_fmt_ts(curr_mtime)}")

        if curr_data.get("in_use"):
            lines.append(f"  ⚠️ Arquivo aberto no Office — conteúdo não lido; mudança detectada por mtime.")
            lines.append("")
            continue

        if "error" in curr_data:
            lines.append(f"  ⚠️ {curr_data['error']}")
            lines.append("")
            continue

        if ftype == "pptx":
            p_count = prev_data.get("slide_count", 0)
            c_count = curr_data.get("slide_count", 0)
            delta   = c_count - p_count
            sign    = "+" if delta >= 0 else ""
            lines.append(f"  Slides: {p_count} → {c_count} ({sign}{delta})")

            prev_titles = {s["n"]: s["title"] for s in prev_data.get("slides", [])}
            curr_titles = {s["n"]: s["title"] for s in curr_data.get("slides", [])}
            new_slides  = [(n, t) for n, t in curr_titles.items() if n not in prev_titles]
            changed     = [(n, t) for n, t in curr_titles.items()
                           if n in prev_titles and prev_titles[n] != t]

            if new_slides:
                lines.append("  Slides novos:")
                for n, t in new_slides[:12]:
                    lines.append(f"    {n}. {t}")
            if changed:
                lines.append("  Títulos alterados:")
                for n, t in changed[:6]:
                    lines.append(f"    {n}. {prev_titles[n][:60]} → {t[:60]}")

        elif ftype == "xlsx":
            prev_set = set(prev_data.get("sheet_names", []))
            curr_set = set(curr_data.get("sheet_names", []))
            added_s   = sorted(curr_set - prev_set)
            removed_s = sorted(prev_set - curr_set)
            if added_s:
                lines.append(f"  Abas novas: {', '.join(added_s)}")
            if removed_s:
                lines.append(f"  Abas removidas: {', '.join(removed_s)}")

            cell_diffs = []
            for sname, sdata in curr_data.get("sheets", {}).items():
                if "cells" not in sdata:
                    continue
                prev_cells = prev_data.get("sheets", {}).get(sname, {}).get("cells", {})
                for ref, cval in sdata["cells"].items():
                    pval = prev_cells.get(ref)
                    if pval != cval:
                        cell_diffs.append(f"    {sname}!{ref}: {pval} → {cval}")
            if cell_diffs:
                lines.append("  Células alteradas:")
                lines.extend(cell_diffs[:8])

        lines.append("")

    if not lines:
        return None

    return f"## 📁 Mudanças no projeto (desde {prev_ts})\n\n" + "\n".join(lines)


# ---------------------------------------------------------------------------
# Meeting materials mode (--meetings)
# ---------------------------------------------------------------------------

def _iter_meeting_files(folder_path, label, cutoff):
    """Itera arquivos relevantes em folder_path e seus subdiretórios imediatos."""
    EXTS = {".pptx", ".xlsx", ".pdf"}
    dirs_to_scan = [folder_path]
    # adiciona subpastas (um nível) — cada sessão/CTM fica numa subpasta datada
    try:
        for entry in os.scandir(folder_path):
            if entry.is_dir():
                dirs_to_scan.append(entry.path)
    except Exception:
        pass

    for dir_path in dirs_to_scan:
        try:
            for entry in os.scandir(dir_path):
                if not entry.is_file():
                    continue
                if entry.name.startswith("~$"):   # arquivo temporário do Office
                    continue
                suffix = Path(entry.name).suffix.lower()
                if suffix not in EXTS:
                    continue
                mtime = entry.stat().st_mtime
                if mtime < cutoff:
                    continue
                yield entry.path, entry.name, suffix, mtime, label
        except Exception:
            continue


def find_meeting_materials(config, days_back=3):
    """
    Varre meeting_folders (+ watch_folders) em busca de PPTX/XLSX modificados
    nos últimos N dias. Escaneia subpastas imediatas (sessões datadas).
    Retorna lista de dicts ordenada por mtime desc.
    """
    cutoff = (datetime.now() - timedelta(days=days_back)).timestamp()
    results = []
    seen = set()

    folders_to_scan = config.get("meeting_folders", []) + config.get("watch_folders", [])

    for folder_cfg in folders_to_scan:
        for fpath, fname, suffix, mtime, label in _iter_meeting_files(
            folder_cfg["path"], folder_cfg["label"], cutoff
        ):
            if fpath in seen:
                continue
            seen.add(fpath)

            item = {
                "path":    fpath,
                "folder":  label,
                "name":    fname,
                "mtime":   mtime,
                "type":    suffix[1:],
                "content": None,
            }
            if suffix == ".pptx":
                item["content"] = extract_pptx_snapshot(fpath)
            results.append(item)

    results.sort(key=lambda x: x["mtime"], reverse=True)
    return results


def build_meetings_report(materials):
    if not materials:
        return "Nenhum arquivo de reunião encontrado nos últimos 3 dias."

    lines = ["## 📋 Materiais de reunião recentes\n"]
    for m in materials:
        lines.append(f"**{m['name']}**")
        lines.append(f"  Pasta: {m['folder']} · Modificado: {_fmt_ts(m['mtime'])}")

        content = m.get("content")
        if content and m["type"] == "pptx":
            if content.get("in_use"):
                lines.append("  ⚠️ Arquivo aberto no Office — slide titles indisponíveis")
            elif "error" not in content:
                count = content.get("slide_count", 0)
                lines.append(f"  {count} slides")
                slides = content.get("slides", [])
                # Mostrar primeiros e últimos títulos (sumário do deck)
                show = slides[:5] + (slides[-2:] if len(slides) > 7 else [])
                seen = set()
                for s in show:
                    if s["n"] not in seen:
                        seen.add(s["n"])
                        lines.append(f"    {s['n']}. {s['title']}")
                if len(slides) > 7:
                    lines.append(f"    … ({len(slides) - 7} slides intermediários)")
        lines.append("")

    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def resolve_snapshot_path(config, base):
    raw = config["snapshot_path"]
    p = Path(raw)
    return str(p) if p.is_absolute() else str(base / raw)


def run_delta(config, snapshot_path, reset=False):
    prev = {} if reset else load_snapshot(snapshot_path)

    curr = {
        "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M"),
        "folders": {},
        "files": {},
    }

    for folder_cfg in config.get("watch_folders", []):
        curr["folders"][folder_cfg["path"]] = scan_folder(folder_cfg["path"])

    for file_cfg in config.get("deep_scan_files", []):
        path  = file_cfg["path"]
        ftype = file_cfg.get("type", "unknown")
        track_cells = file_cfg.get("track_cells")

        # Remove placeholder entries
        if track_cells and "_PLACEHOLDER" in track_cells:
            del track_cells["_PLACEHOLDER"]
            if not track_cells:
                track_cells = None

        try:
            mtime = round(Path(path).stat().st_mtime, 1)
        except FileNotFoundError:
            continue

        if ftype == "pptx":
            data = extract_pptx_snapshot(path)
        elif ftype == "xlsx":
            data = extract_excel_snapshot(path, track_cells)
        else:
            data = {}

        data["mtime"] = mtime
        curr["files"][path] = data

    report = build_delta_report(config, prev, curr)
    save_snapshot(snapshot_path, curr)

    if report:
        _out(report)


def run_meetings(config):
    materials = find_meeting_materials(config, days_back=3)
    report = build_meetings_report(materials)
    _out(report)


def main():
    if len(sys.argv) < 2:
        _out("Uso: python tools/project_scanner.py <config.json> [--reset|--meetings]")
        sys.exit(1)

    config_path = sys.argv[1]
    mode_reset    = "--reset"    in sys.argv
    mode_meetings = "--meetings" in sys.argv

    with open(config_path, encoding="utf-8") as f:
        config = json.load(f)

    base = Path(__file__).parent.parent
    snapshot_path = resolve_snapshot_path(config, base)

    if mode_meetings:
        run_meetings(config)
    else:
        run_delta(config, snapshot_path, reset=mode_reset)


if __name__ == "__main__":
    main()
